# -*- coding: utf-8 -*-
"""
Dify Agent 开发流程 PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==================== 颜色主题 ====================
# 浅色主题 - 温暖的奶油色调
BG_COLOR = RGBColor(0xFA, 0xF7, 0xF2)  # 奶油白
PRIMARY = RGBColor(0x2D, 0x5B, 0x7B)    # 深蓝
SECONDARY = RGBColor(0x4A, 0x90, 0xD9)  # 亮蓝
ACCENT = RGBColor(0xE8, 0x6D, 0x50)     # 珊瑚橙
TEXT_DARK = RGBColor(0x2C, 0x2C, 0x2C)  # 深灰
TEXT_LIGHT = RGBColor(0x6B, 0x6B, 0x6B) # 浅灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD) # 浅蓝背景
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9) # 浅绿背景
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0) # 浅橙背景
LIGHT_PURPLE = RGBColor(0xF3, 0xE5, 0xF5) # 浅紫背景


def set_slide_bg(slide, color):
    """设置幻灯片背景颜色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color):
    """添加圆形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_arrow(slide, left, top, width, height, fill_color):
    """添加箭头"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_points(slide, left, top, width, height, items, font_size=13,
                      color=TEXT_DARK, bullet_color=SECONDARY):
    """添加带项目符号的文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  •  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(6)

    return txBox


def create_title_slide(prs):
    """第1页：封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_bg(slide, BG_COLOR)

    # 装饰圆形
    add_circle(slide, Inches(-0.5), Inches(-0.5), Inches(2), LIGHT_BLUE)
    add_circle(slide, Inches(8.5), Inches(5.5), Inches(2), LIGHT_ORANGE)

    # 标题
    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(7), Inches(1),
                 "Dify Agent 开发流程", font_size=40, color=PRIMARY, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # 副标题
    add_text_box(slide, Inches(1.5), Inches(3.0), Inches(7), Inches(0.6),
                 "从零到一构建你的 AI Agent 应用", font_size=20, color=TEXT_LIGHT,
                 alignment=PP_ALIGN.CENTER)

    # 分割线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.7), Inches(3), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # 日期
    add_text_box(slide, Inches(1.5), Inches(4.2), Inches(7), Inches(0.5),
                 "2026 · 泉 出品", font_size=14, color=TEXT_LIGHT,
                 alignment=PP_ALIGN.CENTER)


def create_overview_slide(prs):
    """第2页：目录概览"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    # 标题
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(4), Inches(0.5),
                 "📋 目录", font_size=28, color=PRIMARY, bold=True)

    items = [
        ("01", "认识 Dify 与 Agent", "了解平台和核心概念"),
        ("02", "环境准备与部署", "SaaS / Docker 自托管"),
        ("03", "创建 Agent 应用", "选择模型、编写 Prompt"),
        ("04", "配置工具 (Tools)", "内置工具 + 自定义 API"),
        ("05", "Agent 策略配置", "推理模式与迭代控制"),
        ("06", "调试与测试", "实时对话与推理追踪"),
        ("07", "发布与集成", "API / 嵌入 / 独立链接"),
    ]

    for i, (num, title, desc) in enumerate(items):
        y = Inches(0.9 + i * 0.6)

        # 编号圆形
        circle = add_circle(slide, Inches(0.7), y, Inches(0.4), SECONDARY)
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = False

        # 标题
        add_text_box(slide, Inches(1.3), y, Inches(3), Inches(0.3),
                     title, font_size=15, color=TEXT_DARK, bold=True)

        # 描述
        add_text_box(slide, Inches(4.5), y, Inches(5), Inches(0.3),
                     desc, font_size=12, color=TEXT_LIGHT)


def create_what_is_dify(prs):
    """第3页：认识 Dify"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(6), Inches(0.5),
                 "01  认识 Dify 与 Agent", font_size=28, color=PRIMARY, bold=True)

    # 左侧 - Dify 介绍
    # 矩形: top=0.9, height=4.5 → bottom=5.4 (安全)
    box1 = add_rounded_rect(slide, Inches(0.5), Inches(0.9), Inches(4.2), Inches(4.5),
                            WHITE, SECONDARY)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(3.8), Inches(0.4),
                 "🚀 Dify 是什么？", font_size=18, color=SECONDARY, bold=True)

    dify_points = [
        "开源的 LLM 应用开发平台",
        "可视化编排 AI 工作流",
        "支持多种大模型接入",
        "提供 RAG、Agent、Chatflow",
        "SaaS 版 + Docker 自托管",
    ]
    # top=1.5, height=3.5 → bottom=5.0 (安全)
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(3.6), Inches(3.5),
                      dify_points, font_size=13)

    # 右侧 - Agent 介绍
    box2 = add_rounded_rect(slide, Inches(5.2), Inches(0.9), Inches(4.2), Inches(4.5),
                            WHITE, ACCENT)
    add_text_box(slide, Inches(5.5), Inches(1.0), Inches(3.8), Inches(0.4),
                 "🤖 Agent 是什么？", font_size=18, color=ACCENT, bold=True)

    agent_points = [
        "能自主推理和决策的 AI",
        "可以调用工具完成任务",
        "支持多轮迭代思考",
        "Thought → Action → Observation",
        "比普通对话更智能",
    ]
    add_bullet_points(slide, Inches(5.5), Inches(1.5), Inches(3.6), Inches(3.5),
                      agent_points, font_size=13)


def create_env_setup(prs):
    """第4页：环境准备"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(6), Inches(0.5),
                 "02  环境准备与部署", font_size=28, color=PRIMARY, bold=True)

    # 方案A - SaaS
    # 矩形: top=0.9, height=4.5 → bottom=5.4 (安全)
    box1 = add_rounded_rect(slide, Inches(0.5), Inches(0.9), Inches(4.2), Inches(4.5),
                            LIGHT_GREEN, None)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(3.8), Inches(0.4),
                 "☁️ 方案 A：SaaS 版（推荐新手）", font_size=15, color=PRIMARY, bold=True)

    saas_steps = [
        "访问 cloud.dify.ai",
        "注册账号（邮箱/GitHub）",
        "直接开始创建应用",
        "免费额度足够学习使用",
        "无需服务器，开箱即用",
    ]
    # top=1.5, height=3.5 → bottom=5.0 (安全)
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(3.6), Inches(3.5),
                      saas_steps, font_size=13)

    # 方案B - Docker
    box2 = add_rounded_rect(slide, Inches(5.2), Inches(0.9), Inches(4.2), Inches(4.5),
                            LIGHT_BLUE, None)
    add_text_box(slide, Inches(5.5), Inches(1.0), Inches(3.8), Inches(0.4),
                 "🐳 方案 B：Docker 自托管", font_size=15, color=PRIMARY, bold=True)

    docker_steps = [
        "安装 Docker + Docker Compose",
        "git clone dify 官方仓库",
        "cd dify/docker",
        "cp .env.example .env",
        "docker compose up -d",
        "访问 localhost/install",
    ]
    add_bullet_points(slide, Inches(5.5), Inches(1.5), Inches(3.6), Inches(3.5),
                      docker_steps, font_size=13)


def create_model_config(prs):
    """第5页：创建Agent应用 & 配置模型"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
                 "03  创建 Agent 应用 & 配置模型", font_size=28, color=PRIMARY, bold=True)

    # 步骤流程
    steps = [
        ("1", "创建应用", "控制台 → 创建应用\n选择「Agent」类型", SECONDARY),
        ("2", "选择模型", "接入 LLM 供应商\n配置 API Key", ACCENT),
        ("3", "编写 Prompt", "定义角色和行为\n设定约束条件", PRIMARY),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        x = Inches(0.5 + i * 3.2)
        y = Inches(1.3)

        # 卡片
        card = add_rounded_rect(slide, x, y, Inches(2.8), Inches(3.5), WHITE, color)

        # 编号
        circle = add_circle(slide, x + Inches(1.05), y + Inches(0.2), Inches(0.6), color)
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 标题
        add_text_box(slide, x + Inches(0.2), y + Inches(1.0), Inches(2.4), Inches(0.4),
                     title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # 描述
        add_text_box(slide, x + Inches(0.2), y + Inches(1.6), Inches(2.4), Inches(1.5),
                     desc, font_size=13, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

        # 箭头（除最后一个）
        if i < 2:
            add_arrow(slide, x + Inches(2.9), y + Inches(1.5), Inches(0.3), Inches(0.3), color)

    # 底部提示
    # top=5.0, height=0.4 → bottom=5.4 (安全)
    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(9), Inches(0.4),
                 "💡 支持的模型：GPT-4、Claude、通义千问、文心一言、DeepSeek 等",
                 font_size=12, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


def create_tools_slide(prs):
    """第6页：配置工具"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(6), Inches(0.5),
                 "04  配置工具 (Tools)", font_size=28, color=PRIMARY, bold=True)

    # 三列工具类型
    tool_types = [
        ("🔧 内置工具", SECONDARY, [
            "网页搜索",
            "计算器",
            "DALL-E 绘图",
            "代码执行器",
            "更多官方工具...",
        ]),
        ("🔌 自定义 API", ACCENT, [
            "编写 OpenAPI Schema",
            "导入 Swagger 文档",
            "对接你的后端服务",
            "支持 GET/POST 请求",
            "可配置认证方式",
        ]),
        ("⚡ 工作流工具", PRIMARY, [
            "将 Workflow 发布为工具",
            "复杂任务拆解",
            "多步骤编排",
            "支持条件分支",
            "可复用逻辑模块",
        ]),
    ]

    for i, (title, color, items) in enumerate(tool_types):
        x = Inches(0.4 + i * 3.2)

        # 标题
        # top=0.9, height=0.5 → bottom=1.4 (安全)
        box = add_rounded_rect(slide, x, Inches(0.9), Inches(2.9), Inches(0.5), color)
        tf = box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 工具列表
        # top=1.5, height=3.5 → bottom=5.0 (安全)
        list_box = add_rounded_rect(slide, x, Inches(1.5), Inches(2.9), Inches(3.5), WHITE, color)
        add_bullet_points(slide, x + Inches(0.2), Inches(1.7), Inches(2.5), Inches(3),
                          items, font_size=12)

    # 底部说明
    # top=5.1, height=0.4 → bottom=5.5 (安全)
    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(9), Inches(0.4),
                 "💡 自定义工具是最强大的扩展方式 —— 任何有 API 的服务都能接入！",
                 font_size=12, color=ACCENT, alignment=PP_ALIGN.CENTER)


def create_strategy_slide(prs):
    """第7页：Agent策略配置"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(6), Inches(0.6),
                 "05  Agent 策略配置", font_size=28, color=PRIMARY, bold=True)

    # 推理模式
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(4), Inches(0.4),
                 "🧠 推理模式选择", font_size=18, color=PRIMARY, bold=True)

    modes = [
        ("Function Calling", "模型原生支持，推荐", "速度快、准确度高", LIGHT_GREEN),
        ("ReAct", "通用推理框架", "兼容性好，适合所有模型", LIGHT_BLUE),
    ]

    for i, (name, desc, detail, bg) in enumerate(modes):
        x = Inches(0.5 + i * 4.5)
        box = add_rounded_rect(slide, x, Inches(1.8), Inches(4), Inches(1.5), bg, None)
        add_text_box(slide, x + Inches(0.2), Inches(1.9), Inches(3.6), Inches(0.4),
                     name, font_size=16, color=PRIMARY, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(2.3), Inches(3.6), Inches(0.3),
                     desc, font_size=13, color=TEXT_DARK)
        add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(3.6), Inches(0.3),
                     detail, font_size=12, color=TEXT_LIGHT)

    # 关键参数
    # top=3.5, height=0.4 → bottom=3.9 (安全)
    add_text_box(slide, Inches(0.5), Inches(3.5), Inches(4), Inches(0.4),
                 "⚙️ 关键参数", font_size=18, color=PRIMARY, bold=True)

    params = [
        ("最大迭代轮数", "控制 Agent 思考深度，建议 3-5 轮", "🔄"),
        ("Temperature", "0 = 精确，1 = 创意，Agent 建议 0.7", "🌡️"),
        ("创意/精确度", "滑块调节，影响回答风格", "🎯"),
    ]

    for i, (name, desc, emoji) in enumerate(params):
        # y=4.0/4.4/4.8, height=0.4 → max bottom=5.2 (安全)
        y = Inches(4.0 + i * 0.45)
        add_text_box(slide, Inches(0.8), y, Inches(2), Inches(0.35),
                     f"{emoji} {name}", font_size=13, color=TEXT_DARK, bold=True)
        add_text_box(slide, Inches(3), y, Inches(6.5), Inches(0.35),
                     desc, font_size=13, color=TEXT_LIGHT)


def create_debug_slide(prs):
    """第8页：调试与测试"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(6), Inches(0.5),
                 "06  调试与测试", font_size=28, color=PRIMARY, bold=True)

    # 推理过程可视化
    add_text_box(slide, Inches(0.5), Inches(0.9), Inches(9), Inches(0.4),
                 "🔍 Agent 推理过程（Thought → Action → Observation）",
                 font_size=16, color=PRIMARY, bold=True)

    # 三个步骤卡片
    thought_steps = [
        ("💭 Thought", "Agent 的思考过程\n分析问题，制定计划\n决定需要什么信息", SECONDARY),
        ("⚡ Action", "执行具体操作\n调用工具/API\n搜索、计算、查询", ACCENT),
        ("👁️ Observation", "观察执行结果\n获取返回数据\n判断是否完成", PRIMARY),
    ]

    for i, (title, desc, color) in enumerate(thought_steps):
        x = Inches(0.5 + i * 3.2)
        card = add_rounded_rect(slide, x, Inches(1.4), Inches(2.8), Inches(2.2), WHITE, color)

        add_text_box(slide, x + Inches(0.2), Inches(1.5), Inches(2.4), Inches(0.4),
                     title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(2.4), Inches(1.2),
                     desc, font_size=12, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)

        if i < 2:
            add_arrow(slide, x + Inches(2.9), Inches(2.3), Inches(0.3), Inches(0.25), color)

    # 调试技巧
    add_text_box(slide, Inches(0.5), Inches(3.9), Inches(4), Inches(0.4),
                 "💡 调试技巧", font_size=16, color=PRIMARY, bold=True)

    tips = [
        "查看推理日志，定位问题环节",
        "逐步增加迭代轮数，观察行为变化",
        "用简单问题先验证基本功能",
        "检查工具返回数据是否符合预期",
    ]
    add_bullet_points(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(1.2),
                      tips, font_size=12)


def create_publish_slide(prs):
    """第9页：发布与集成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(6), Inches(0.5),
                 "07  发布与集成", font_size=28, color=PRIMARY, bold=True)

    # 三种发布方式
    publish_methods = [
        ("🌐 API 发布", SECONDARY, [
            "生成 API Key",
            "RESTful 接口调用",
            "支持流式响应",
            "可对接任何后端",
            "适合生产环境",
        ]),
        ("📱 嵌入集成", ACCENT, [
            "iframe 嵌入网页",
            "JS SDK 集成",
            "自定义 UI 样式",
            "适合产品集成",
            "一行代码接入",
        ]),
        ("🔗 独立链接", PRIMARY, [
            "生成访问链接",
            "可设置访问密码",
            "分享给团队使用",
            "适合内部工具",
            "无需开发",
        ]),
    ]

    for i, (title, color, items) in enumerate(publish_methods):
        x = Inches(0.4 + i * 3.2)

        # 标题
        # top=0.9, height=0.5 → bottom=1.4 (安全)
        box = add_rounded_rect(slide, x, Inches(0.9), Inches(2.9), Inches(0.5), color)
        tf = box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 列表
        # top=1.5, height=3.3 → bottom=4.8 (安全)
        list_box = add_rounded_rect(slide, x, Inches(1.5), Inches(2.9), Inches(3.3), WHITE, color)
        add_bullet_points(slide, x + Inches(0.2), Inches(1.7), Inches(2.5), Inches(2.8),
                          items, font_size=12)

    # API 示例代码
    # top=5.0, height=0.4 → bottom=5.4 (安全)
    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(9), Inches(0.4),
                 '💻 API 调用示例：POST /v1/chat-messages  { "query": "你好", "user": "user-123" }',
                 font_size=11, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


def create_summary_slide(prs):
    """第10页：总结"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    # 装饰
    add_circle(slide, Inches(-0.5), Inches(-0.5), Inches(2), LIGHT_BLUE)
    add_circle(slide, Inches(8.5), Inches(4.5), Inches(1.5), LIGHT_ORANGE)

    add_text_box(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.6),
                 "✨ 开发流程总结", font_size=30, color=PRIMARY, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # 流程图
    flow_steps = [
        ("准备环境", "注册/部署"),
        ("创建应用", "选模型"),
        ("写 Prompt", "定角色"),
        ("配置工具", "接 API"),
        ("策略调优", "设参数"),
        ("调试测试", "查日志"),
        ("发布上线", "集成产品"),
    ]

    for i, (step, detail) in enumerate(flow_steps):
        x = Inches(0.3 + i * 1.35)
        y = Inches(1.2)

        # 步骤卡片
        color = [SECONDARY, ACCENT, PRIMARY, SECONDARY, ACCENT, PRIMARY, SECONDARY][i]
        box = add_rounded_rect(slide, x, y, Inches(1.2), Inches(1.0), WHITE, color)

        add_text_box(slide, x + Inches(0.05), y + Inches(0.1), Inches(1.1), Inches(0.35),
                     step, font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + Inches(0.05), y + Inches(0.5), Inches(1.1), Inches(0.35),
                     detail, font_size=9, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

        if i < 6:
            add_arrow(slide, x + Inches(1.25), y + Inches(0.35), Inches(0.1), Inches(0.15), color)

    # 学习资源
    add_text_box(slide, Inches(1), Inches(2.5), Inches(8), Inches(0.4),
                 "📚 学习资源", font_size=18, color=PRIMARY, bold=True,
                 alignment=PP_ALIGN.CENTER)

    resources = [
        "官方文档：docs.dify.ai",
        "GitHub：github.com/langgenius/dify",
        "B站搜索「Dify Agent 教程」",
        "知乎搜索「Dify 实战」",
    ]

    for i, res in enumerate(resources):
        x = Inches(1.5 + (i % 2) * 4)
        y = Inches(3.0 + (i // 2) * 0.45)
        add_text_box(slide, x, y, Inches(3.5), Inches(0.35),
                     f"  •  {res}", font_size=13, color=TEXT_DARK)

    # 结束语
    add_text_box(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.5),
                 "🚀 开始你的 Agent 开发之旅吧！有问题随时找我~",
                 font_size=16, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


def main():
    """主函数"""
    prs = Presentation()

    # 设置 16:9 比例
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 创建所有幻灯片
    create_title_slide(prs)       # 封面
    create_overview_slide(prs)    # 目录
    create_what_is_dify(prs)      # 认识 Dify
    create_env_setup(prs)         # 环境准备
    create_model_config(prs)      # 创建应用 & 模型配置
    create_tools_slide(prs)       # 配置工具
    create_strategy_slide(prs)    # Agent 策略
    create_debug_slide(prs)       # 调试测试
    create_publish_slide(prs)     # 发布集成
    create_summary_slide(prs)     # 总结

    # 保存文件
    output_path = "D:/pycharm/Person-Practice/PPT/Dify_Agent开发流程.pptx"
    prs.save(output_path)
    print(f"PPT done: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
